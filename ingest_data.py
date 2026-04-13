import asyncio
import csv
import os
from pathlib import Path

import psycopg
from docx import Document
from dotenv import load_dotenv
from pgvector.psycopg import register_vector_async
from psycopg.types.json import Json
from pptx import Presentation
from sentence_transformers import SentenceTransformer

load_dotenv()

# --- Configuration ---
DB_HOST = os.getenv("DB_HOST", "localhost")
DB_NAME = os.getenv("DB_NAME")
DB_USER = os.getenv("DB_USER")
DB_PASSWORD = os.getenv("DB_PASSWORD")

EMBEDDING_MODEL_NAME = "sentence-transformers/all-MiniLM-L6-v2"
BATCH_SIZE = 50

BASE_DIR = Path(__file__).resolve().parent

# Each entry: (source_name, documents_dir, metadata_csv)
SOURCES = [
    (
        "dropbox",
        BASE_DIR / "data" / "dropbox" / "documents",
        BASE_DIR / "data" / "dropbox" / "documents" / "metadata.csv",
    ),
    (
        "drive",
        BASE_DIR / "data" / "drive" / "documents",
        BASE_DIR / "data" / "drive" / "documents" / "metadata.csv",
    ),
]


# --- Text extraction ---

def extract_text_txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace").strip()


def extract_text_docx(path: Path) -> str:
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_pptx(path: Path) -> str:
    prs = Presentation(str(path))
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def extract_text(path: Path) -> str | None:
    ext = path.suffix.lower()
    if ext == ".txt":
        return extract_text_txt(path)
    if ext == ".docx":
        return extract_text_docx(path)
    if ext == ".pptx":
        return extract_text_pptx(path)
    return None


# --- Metadata loading ---

def load_metadata_index(csv_path: Path) -> dict[str, dict]:
    """Return a dict mapping downloaded_name -> metadata row."""
    if not csv_path.exists():
        return {}
    with csv_path.open(newline="") as f:
        return {row["downloaded_name"]: row for row in csv.DictReader(f)}


# --- Main ingestion ---

async def ingest_data():
    print("Starting data ingestion...")

    print(f"Loading embedding model: '{EMBEDDING_MODEL_NAME}'")
    model = SentenceTransformer(EMBEDDING_MODEL_NAME)

    conn = None
    try:
        conn_string = f"host={DB_HOST} dbname={DB_NAME} user={DB_USER} password={DB_PASSWORD}"
        conn = await psycopg.AsyncConnection.connect(conn_string)
        await register_vector_async(conn)
        print("Connected to database.")

        # Load already-ingested source IDs to avoid duplicates
        async with conn.cursor() as cur:
            await cur.execute("SELECT metadata->>'source_id' FROM documents WHERE metadata->>'source_id' IS NOT NULL")
            rows = await cur.fetchall()
        ingested_ids = {row[0] for row in rows}
        print(f"Already ingested: {len(ingested_ids)} document(s).")

        total = 0
        for source_name, docs_dir, metadata_csv in SOURCES:
            if not docs_dir.exists():
                print(f"No documents directory for {source_name}, skipping.")
                continue

            meta_index = load_metadata_index(metadata_csv)
            files = [p for p in docs_dir.iterdir() if p.is_file() and not p.name.startswith("._") and p.suffix.lower() in {".txt", ".docx", ".pptx"}]

            if not files:
                print(f"No files found in {docs_dir}")
                continue

            print(f"\n{source_name}: {len(files)} file(s) found.")

            to_ingest = []
            for file_path in files:
                meta = meta_index.get(file_path.name, {})
                source_id = meta.get("id", "")

                if source_id and source_id in ingested_ids:
                    print(f"  Skipping (already ingested): {file_path.name}")
                    continue

                try:
                    content = extract_text(file_path)
                except Exception as e:
                    print(f"  Skipping (extraction error): {file_path.name} — {e}")
                    continue
                if content:
                    content = content.replace("\x00", "")
                if not content:
                    print(f"  Skipping (no extractable text): {file_path.name}")
                    continue

                to_ingest.append((file_path.name, content, {
                    "source": source_name,
                    "source_id": source_id,
                    "original_name": meta.get("name", file_path.name),
                    "path": meta.get("path", ""),
                }))

            if not to_ingest:
                print(f"  Nothing new to ingest from {source_name}.")
                continue

            print(f"  {len(to_ingest)} file(s) to ingest, processing in batches of {BATCH_SIZE}...")
            for batch_start in range(0, len(to_ingest), BATCH_SIZE):
                batch = to_ingest[batch_start : batch_start + BATCH_SIZE]
                batch_num = batch_start // BATCH_SIZE + 1
                total_batches = (len(to_ingest) + BATCH_SIZE - 1) // BATCH_SIZE
                print(f"  Batch {batch_num}/{total_batches} — embedding {len(batch)} file(s)...")

                contents = [item[1] for item in batch]
                embeddings = model.encode(contents, show_progress_bar=False)

                async with conn.cursor() as cur:
                    for i, (filename, content, metadata) in enumerate(batch):
                        await cur.execute(
                            "INSERT INTO documents (content, metadata, embedding) VALUES (%s, %s, %s)",
                            (content, Json(metadata), embeddings[i]),
                        )
                    await conn.commit()
                print(f"  Batch {batch_num}/{total_batches} committed.")

            total += len(to_ingest)

        print(f"\nIngestion complete. {total} new document(s) added.")

    except Exception as e:
        print(f"Error during ingestion: {e}")
        raise
    finally:
        if conn:
            await conn.close()


if __name__ == "__main__":
    asyncio.run(ingest_data())
